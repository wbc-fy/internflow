from io import BytesIO
from pathlib import Path

from fastapi import UploadFile
from pptx import Presentation


MAX_FILE_SIZE = 5 * 1024 * 1024


async def extract_text_from_pptx_upload(file: UploadFile) -> str:
    if not file.filename:
        raise ValueError("Filename is missing.")

    file_ext = Path(file.filename).suffix.lower()

    if file_ext != ".pptx":
        raise ValueError("Only .pptx files are supported now.")

    content = await file.read()

    if len(content) > MAX_FILE_SIZE:
        raise ValueError("File is too large. Max size is 5MB.")

    return extract_text_from_pptx(content)


def extract_text_from_pptx(content: bytes) -> str:
    presentation = Presentation(BytesIO(content))
    texts = []

    for slide_index, slide in enumerate(presentation.slides, start=1):
        slide_texts = []

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text.strip())

        if slide_texts:
            texts.append(f"Slide {slide_index}:\n" + "\n".join(slide_texts))

    extracted_text = "\n\n".join(texts).strip()

    if not extracted_text:
        raise ValueError(
            "No text found in this PPTX. It may contain only images or screenshots."
        )

    return extracted_text